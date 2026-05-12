import importlib.util
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = "data_prep_progress.pptx"
DB_CANDIDATES = [Path("macys.db"), Path("data/macys.db")]
SKILL_PATH = Path("skills/audience-segment-builder/segment.py")

CREAM = RGBColor(0xF8, 0xF4, 0xEC)
TEAL = RGBColor(0x0B, 0x7B, 0x8A)
CHARCOAL = RGBColor(0x2D, 0x2D, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEAL = RGBColor(0xE3, 0xEE, 0xF0)
LIGHTER_TEAL = RGBColor(0xF1, 0xF7, 0xF8)
DONE_GREEN = RGBColor(0x2E, 0x8B, 0x57)
IN_PROGRESS_AMBER = RGBColor(0xD4, 0x9A, 0x12)
CODE_BG = RGBColor(0x1F, 0x2E, 0x33)

TITLE_FONT = "Georgia"
BODY_FONT = "Calibri"
MONO_FONT = "Consolas"

SPEAKER_NOTES = [
    # Slide 1: Title
    "Hi everyone, we are Group 4. Today we are sharing our progress on the data "
    "preparation for our GenAI co-worker for Macy's marketing campaign execution. "
    "Quick five minute walkthrough of what we built and where we are headed.",
    # Slide 2: Workflow → Skills → Data
    "This is the core of our project mapped out in three parts. On the left, the "
    "10 step workflow we documented in M1. In the middle, the four GenAI skills "
    "we are building, each one targeting a specific bottleneck step. On the "
    "right, the data each skill needs to do its job. The whole project is built "
    "around making sure these three columns line up.",
    # Slide 3: How We Prepared the Data
    "We didn't just download a dataset, we built ours from the ground up. "
    "Researched real retail operations, designed a 7 table schema that mirrors a "
    "real Macy's data warehouse, then simulated with Faker plus LLM. We even "
    "injected the same kind of mess you would find in real retail data. The "
    "result is a 480 thousand row database we can query from any of our skills.",
    # Slide 4: Skill 1 Live
    "This is our first working skill, the Audience Segment Builder. Merna asks a "
    "question, the skill queries our real database, and returns three ranked "
    "segment options. Beauty Loyalists, Gift Givers, First Time Buyers, with "
    "real customer counts and average spend. Eighteen passing tests. Happy to "
    "demo it live if there is time.",
    # Slide 5: What is Next
    "Skill 1 is done and ready. Three more in the pipeline this week, each one "
    "targeting a specific workflow bottleneck we identified in M1. Final 10 "
    "minute demo video on May 7 with all four skills running live. Repo link in "
    "the footer if you want to look at the code.",
]


def find_db():
    for p in DB_CANDIDATES:
        if p.exists():
            return p
    raise FileNotFoundError(f"macys.db not found in any of: {DB_CANDIDATES}")


def load_skill_module(path: Path):
    spec = importlib.util.spec_from_file_location("audience_segment_builder", path)
    if spec is None or spec.loader is None:
        raise ImportError(f"Could not load skill module at {path}")
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


db_path = find_db()
skill = load_skill_module(SKILL_PATH)
live_segments = skill.build_segments("Mother's Day Beauty Event", db_path=db_path)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
blank = prs.slide_layouts[6]


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = CREAM


def textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    return tf


def style_run(run, *, font=BODY_FONT, size=14, color=CHARCOAL, bold=False, italic=False):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def add_para(
    tf,
    text,
    *,
    font=BODY_FONT,
    size=14,
    color=CHARCOAL,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    space_after=None,
    space_before=None,
):
    if len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    if space_after is not None:
        p.space_after = Pt(space_after)
    if space_before is not None:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    style_run(run, font=font, size=size, color=color, bold=bold, italic=italic)
    return p


def teal_rule(slide, x, y, w, h=0.04):
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rule.line.fill.background()
    rule.fill.solid()
    rule.fill.fore_color.rgb = TEAL


def style_cell(cell, *, header=False, zebra_idx=0, font=BODY_FONT, size=12, align=PP_ALIGN.LEFT, bold=False):
    cell.fill.solid()
    if header:
        cell.fill.fore_color.rgb = TEAL
    else:
        cell.fill.fore_color.rgb = CREAM if zebra_idx % 2 == 1 else LIGHTER_TEAL
    cell.margin_left = Inches(0.1)
    cell.margin_right = Inches(0.1)
    cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.05)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = cell.text_frame.paragraphs[0]
    p.alignment = align
    for run in p.runs:
        if header:
            style_run(run, font=BODY_FONT, size=size, color=WHITE, bold=True)
        else:
            style_run(run, font=font, size=size, color=CHARCOAL, bold=bold)


def numbered_badge(slide, x, y, diameter, number):
    """Draw a teal filled circle with a white number inside."""
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(diameter), Inches(diameter))
    badge.fill.solid()
    badge.fill.fore_color.rgb = TEAL
    badge.line.fill.background()
    tf = badge.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    style_run(run, font=TITLE_FONT, size=14, color=WHITE, bold=True)


# ---------------- SLIDE 1: Title ----------------
s1 = prs.slides.add_slide(blank)
add_bg(s1)

tf = textbox(s1, 0.5, 1.6, 12.333, 0.5)
add_para(tf, "MILESTONE 2 PROGRESS UPDATE", size=16, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

teal_rule(s1, 6.17, 2.25, 1.0)

tf = textbox(s1, 0.5, 2.6, 12.333, 1.4)
add_para(
    tf,
    "Macy's Marketing Campaign Execution",
    font=TITLE_FONT,
    size=44,
    color=CHARCOAL,
    bold=True,
    align=PP_ALIGN.CENTER,
)

tf = textbox(s1, 0.5, 4.05, 12.333, 0.7)
add_para(
    tf, "GenAI Co-worker for Omnichannel Retail Operations", size=22, color=TEAL, italic=True, align=PP_ALIGN.CENTER
)

tf = textbox(s1, 0.5, 5.6, 12.333, 0.5)
add_para(tf, "Group 4: Merna Saad, Abdullah, Shankar", size=16, color=CHARCOAL, align=PP_ALIGN.CENTER)

tf = textbox(s1, 0.5, 6.1, 12.333, 0.5)
add_para(tf, "MGT 449, Spring 2026", size=14, color=CHARCOAL, italic=True, align=PP_ALIGN.CENTER)


# ---------------- SLIDE 2: Workflow → Skills → Data ----------------
s2 = prs.slides.add_slide(blank)
add_bg(s2)

tf = textbox(s2, 0.5, 0.35, 12.333, 0.7)
add_para(
    tf, "From Workflow to Skills to Data", font=TITLE_FONT, size=30, color=CHARCOAL, bold=True, align=PP_ALIGN.CENTER
)

teal_rule(s2, 5.67, 1.05, 2.0)

col_y = 1.35
col_h = 5.2
col_w = 4.0
col_xs = [0.5, 4.67, 8.83]

tf = textbox(s2, col_xs[0], col_y, col_w, 0.5)
add_para(tf, "10 STEP WORKFLOW", size=15, color=TEAL, bold=True)

steps = [
    "1. Briefing",
    "2. Segmentation",
    "3. SKU Selection",
    "4. Creative Production",
    "5. Layout Assembly",
    "6. Approval",
    "7. Localization",
    "8. Activation",
    "9. Monitoring",
    "10. Reporting",
]
tf = textbox(s2, col_xs[0], col_y + 0.5, col_w, col_h - 0.5)
for step in steps:
    add_para(tf, step, size=13, color=CHARCOAL, space_after=4)

tf = textbox(s2, col_xs[1], col_y, col_w, 0.5)
add_para(tf, "4 SKILLS", size=15, color=TEAL, bold=True)

skills_list = [
    ("SKILL 1: Audience Segment Builder", "Targets: Step 2"),
    ("SKILL 2: DAM Asset Finder", "Targets: Steps 4 to 5"),
    ("SKILL 3: Localization Generator", "Targets: Step 7"),
    ("SKILL 4: Campaign Performance Analyzer", "Targets: Step 9"),
]
tf = textbox(s2, col_xs[1], col_y + 0.5, col_w, col_h - 0.5)
for name, target in skills_list:
    add_para(tf, name, size=13, color=CHARCOAL, bold=True, space_after=2)
    add_para(tf, target, size=12, color=TEAL, italic=True, space_after=10)

tf = textbox(s2, col_xs[2], col_y, col_w, 0.5)
add_para(tf, "DATA NEEDED", size=15, color=TEAL, bold=True)

data_items = [
    "Customer profiles, transactions, loyalty data",
    "DAM asset metadata, tags, rights, resolution",
    "Master ad spec, regional pricing, inventory",
    "Campaign performance by channel and segment",
]
tf = textbox(s2, col_xs[2], col_y + 0.5, col_w, col_h - 0.5)
for item in data_items:
    add_para(tf, "•  " + item, size=13, color=CHARCOAL, space_after=14)

tf = textbox(s2, 0.5, 6.75, 12.333, 0.5)
add_para(
    tf,
    "Each skill solves one workflow bottleneck. Each one needs specific data.",
    size=14,
    color=TEAL,
    italic=True,
    align=PP_ALIGN.CENTER,
)


# ---------------- SLIDE 3: How We Prepared the Data ----------------
s3 = prs.slides.add_slide(blank)
add_bg(s3)

tf = textbox(s3, 0.5, 0.35, 12.333, 0.6)
add_para(tf, "How We Prepared the Data", font=TITLE_FONT, size=30, color=CHARCOAL, bold=True, align=PP_ALIGN.CENTER)

tf = textbox(s3, 0.5, 0.95, 12.333, 0.45)
add_para(
    tf,
    "From raw industry research to working database",
    size=15,
    color=TEAL,
    italic=True,
    align=PP_ALIGN.CENTER,
)

# LEFT: 5-step process
tf = textbox(s3, 0.5, 1.55, 6.0, 0.4)
add_para(tf, "OUR PROCESS", size=14, color=TEAL, bold=True)

process_steps = [
    (
        "RESEARCHED MACY'S OPERATIONS",
        "Pulled industry data on retail marketing cycles, ROI benchmarks, and DAM challenges to ground the simulation in reality.",
    ),
    (
        "DESIGNED THE SCHEMA",
        "7 normalized tables to mirror a real Macy's marketing data warehouse: customers, transactions, SKUs, regional pricing, DAM assets, campaigns, performance.",
    ),
    (
        "SIMULATED WITH FAKER + LLM",
        "Faker for structural data, LLM personas for realistic copy. Postgres-compatible SQLite for portability.",
    ),
    (
        "INJECTED REALISTIC MESS",
        "2% null emails, 30% degraded DAM assets, expired rights, mislabeled tags. Real retail data is never clean.",
    ),
    (
        "PULLED REAL PRODUCT PHOTOS",
        "50 images via the Unsplash API populate the DAM tables, anchoring the visual side of the demo.",
    ),
]

step_origin_y = 2.0
step_pitch = 0.95
text_x = 1.15
text_w = 5.35

for i, (title, body) in enumerate(process_steps):
    y = step_origin_y + i * step_pitch
    numbered_badge(s3, 0.5, y, 0.42, i + 1)
    tf = textbox(s3, text_x, y - 0.04, text_w, 0.32)
    add_para(tf, title, size=12, color=CHARCOAL, bold=True)
    tf = textbox(s3, text_x, y + 0.28, text_w, 0.6)
    add_para(tf, body, size=10.5, color=CHARCOAL)

# RIGHT: produced summary table
tf = textbox(s3, 7.0, 1.55, 5.83, 0.4)
add_para(tf, "WHAT WE PRODUCED", size=14, color=TEAL, bold=True)

table_summary = [
    ("customers", "50,000"),
    ("transactions", "417,180"),
    ("sku_catalog", "2,000"),
    ("regional_pricing", "20,000"),
    ("dam_assets", "5,000"),
    ("campaigns", "12"),
    ("campaign_performance", "2,920"),
]
n = len(table_summary) + 1
right_tbl = s3.shapes.add_table(n, 2, Inches(7.0), Inches(2.0), Inches(5.83), Inches(3.6)).table
right_tbl.columns[0].width = Inches(3.7)
right_tbl.columns[1].width = Inches(2.13)

right_tbl.cell(0, 0).text = "Table"
right_tbl.cell(0, 1).text = "Rows"
style_cell(right_tbl.cell(0, 0), header=True, size=12)
style_cell(right_tbl.cell(0, 1), header=True, size=12, align=PP_ALIGN.RIGHT)
for r, (name, rows_str) in enumerate(table_summary, start=1):
    right_tbl.cell(r, 0).text = name
    right_tbl.cell(r, 1).text = rows_str
    style_cell(right_tbl.cell(r, 0), zebra_idx=r, font=MONO_FONT, size=12, bold=True)
    style_cell(right_tbl.cell(r, 1), zebra_idx=r, font=MONO_FONT, size=12, align=PP_ALIGN.RIGHT)

tf = textbox(s3, 7.0, 5.75, 5.83, 1.0)
add_para(
    tf,
    "All 480K+ rows seeded for reproducibility (seed=42).",
    size=11,
    color=CHARCOAL,
    italic=True,
    space_after=2,
)
add_para(
    tf,
    "Database lives in data/macys.db, queryable from any of our skills.",
    size=11,
    color=CHARCOAL,
    italic=True,
)


# ---------------- SLIDE 4: Skill 1 Live ----------------
s4 = prs.slides.add_slide(blank)
add_bg(s4)

tf = textbox(s4, 0.5, 0.35, 12.333, 0.6)
add_para(
    tf,
    "Skill 1 Live: Audience Segment Builder",
    font=TITLE_FONT,
    size=28,
    color=CHARCOAL,
    bold=True,
    align=PP_ALIGN.CENTER,
)

tf = textbox(s4, 0.5, 0.95, 12.333, 0.45)
add_para(
    tf,
    "From Merna's question to a real segment in 90 seconds",
    size=15,
    color=TEAL,
    italic=True,
    align=PP_ALIGN.CENTER,
)

content_x = 1.0
content_w = 11.33

# Section 1: Merna's input, terminal style command
y = 1.55
tf = textbox(s4, content_x, y, content_w, 0.32)
add_para(tf, "MERNA'S INPUT", size=12, color=TEAL, bold=True)

cmd_box = s4.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(content_x), Inches(y + 0.34), Inches(content_w), Inches(0.55)
)
cmd_box.adjustments[0] = 0.15
cmd_box.fill.solid()
cmd_box.fill.fore_color.rgb = CODE_BG
cmd_box.line.fill.background()
cmd_box.text_frame.text = ""

tf = textbox(s4, content_x + 0.2, y + 0.4, content_w - 0.4, 0.45)
add_para(
    tf,
    '$ python skills/audience-segment-builder/segment.py "Mother\'s Day Beauty Event"',
    font=MONO_FONT,
    size=13,
    color=WHITE,
)

# Section 2: What the skill does
y = 2.55
tf = textbox(s4, content_x, y, content_w, 0.32)
add_para(tf, "WHAT THE SKILL DOES", size=12, color=TEAL, bold=True)

tf = textbox(s4, content_x, y + 0.34, content_w, 1.0)
for line in [
    "Connects to data/macys.db",
    "Joins customers, transactions, and sku_catalog tables",
    "Returns 3 ranked segment options with real demographics",
]:
    add_para(tf, "•  " + line, size=14, color=CHARCOAL, space_after=4)

# Section 3: Real output table
y = 4.0
tf = textbox(s4, content_x, y, content_w, 0.32)
add_para(tf, "REAL OUTPUT FROM OUR DATABASE", size=12, color=TEAL, bold=True)

# Pull live values from the skill
seg_rows = []
for s in live_segments:
    seg_rows.append(
        (
            s["name"],
            f"{s['customer_count']:,}",
            f"${s['avg_total_spend']:,.0f}",
        )
    )

out_headers = ["Segment", "Customers", "Avg Spend"]
n = len(seg_rows) + 1
out_tbl = s4.shapes.add_table(n, 3, Inches(content_x), Inches(y + 0.36), Inches(content_w), Inches(1.85)).table
out_tbl.columns[0].width = Inches(5.5)
out_tbl.columns[1].width = Inches(2.83)
out_tbl.columns[2].width = Inches(3.0)

for c, h in enumerate(out_headers):
    out_tbl.cell(0, c).text = h
    align = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
    style_cell(out_tbl.cell(0, c), header=True, size=13, align=align)

for r, row in enumerate(seg_rows, start=1):
    for c, val in enumerate(row):
        out_tbl.cell(r, c).text = val
        align = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
        font = BODY_FONT if c == 0 else MONO_FONT
        bold = c == 0
        style_cell(out_tbl.cell(r, c), zebra_idx=r, font=font, size=13, align=align, bold=bold)

tf = textbox(s4, 0.5, 6.7, 12.333, 0.5)
add_para(
    tf,
    "4 hours of manual Excel work becomes 90 seconds. 18 passing tests. Live demo available.",
    size=13,
    color=TEAL,
    italic=True,
    align=PP_ALIGN.CENTER,
)


# ---------------- SLIDE 5: What is next ----------------
s5 = prs.slides.add_slide(blank)
add_bg(s5)

tf = textbox(s5, 0.5, 0.35, 12.333, 0.7)
add_para(tf, "What is Next", font=TITLE_FONT, size=30, color=CHARCOAL, bold=True, align=PP_ALIGN.CENTER)

teal_rule(s5, 6.17, 1.05, 1.0)

col_y = 1.55
col_w = 5.8
left_x = 0.7
right_x = 6.83

tf = textbox(s5, left_x, col_y, col_w, 0.5)
add_para(tf, "BUILDING THIS WEEK", size=16, color=TEAL, bold=True)

skill_status = [
    ("Skill 1: Audience Segment Builder", "DONE, demo today", "done"),
    ("Skill 2: DAM Asset Finder", "in progress", "in_progress"),
    ("Skill 3: Localization Generator", "in progress", "in_progress"),
    ("Skill 4: Campaign Performance Analyzer", "in progress", "in_progress"),
]
tf = textbox(s5, left_x, col_y + 0.6, col_w, 3.5)
first = True
for name, status_text, status_kind in skill_status:
    if first:
        p = tf.paragraphs[0]
        first = False
    else:
        p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(8)
    icon = p.add_run()
    if status_kind == "done":
        icon.text = "✓  "
        style_run(icon, font=BODY_FONT, size=15, color=DONE_GREEN, bold=True)
    else:
        icon.text = "○  "
        style_run(icon, font=BODY_FONT, size=15, color=IN_PROGRESS_AMBER, bold=True)
    name_run = p.add_run()
    name_run.text = name + "  "
    style_run(name_run, font=BODY_FONT, size=14, color=CHARCOAL, bold=True)
    status_run = p.add_run()
    status_run.text = f"({status_text})"
    color = DONE_GREEN if status_kind == "done" else IN_PROGRESS_AMBER
    style_run(status_run, font=BODY_FONT, size=12, color=color, italic=True)

tf = textbox(s5, left_x, col_y + 4.4, col_w, 0.6)
add_para(
    tf,
    "README that maps every workflow task to its skill, data, and code.",
    size=12,
    color=CHARCOAL,
    italic=True,
)

tf = textbox(s5, right_x, col_y, col_w, 0.5)
add_para(tf, "DEMO ON MAY 7", size=16, color=TEAL, bold=True)

demo = [
    "10-minute video walkthrough",
    "All 4 skills running live against the database",
    "Real outputs for the Mother's Day campaign scenario",
    "Team narration without scripts",
]
tf = textbox(s5, right_x, col_y + 0.6, col_w, 4.5)
for item in demo:
    add_para(tf, "•  " + item, size=15, color=CHARCOAL, space_after=12)

tf = textbox(s5, 0.5, 6.7, 12.333, 0.5)
add_para(
    tf,
    "Repository: github.com/rsm-genai-2026/mgt449-milestone02-Group-4",
    size=13,
    color=TEAL,
    italic=True,
    align=PP_ALIGN.CENTER,
)


for slide, note in zip(prs.slides, SPEAKER_NOTES):
    slide.notes_slide.notes_text_frame.text = note

prs.save(OUT)
print(f"Saved {len(prs.slides)} slides to: {OUT}")
print(f"Speaker notes added to {len(SPEAKER_NOTES)} slides.")
print(f"DB used: {db_path}")
print(f"Skill module loaded from: {SKILL_PATH}")
print()
print("Slide 4 — segment table populated live from skills/audience-segment-builder/segment.py:")
for s in live_segments:
    print(f"  {s['name']:<20s}  count={s['customer_count']:>7,}  avg_spend=${s['avg_total_spend']:>6,.0f}")
